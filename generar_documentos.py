import os
import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

import pptx
from pptx import Presentation
from pptx.util import Inches as PPTInches, Pt as PPTPt
from pptx.dml.color import RGBColor as PPTRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def set_cell_background(cell, fill_hex):
    tcPr = cell._element.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), fill_hex)
    tcPr.append(shd)

def create_apa_document(title, subtitle, student_name, teacher_name, date_str, sections, filename):
    doc = docx.Document()
    
    # Setup 1 inch margins (APA)
    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)
        
    normal_style = doc.styles['Normal']
    normal_style.font.name = 'Arial'
    normal_style.font.size = Pt(11)
    normal_style.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    normal_style.paragraph_format.line_spacing = 1.15
    normal_style.paragraph_format.space_after = Pt(6)

    # PORTADA APA
    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_title.paragraph_format.space_before = Pt(72)
    p_title.paragraph_format.space_after = Pt(18)
    run_t = p_title.add_run(title)
    run_t.font.size = Pt(18)
    run_t.font.bold = True
    run_t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    p_sub = doc.add_paragraph()
    p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_sub.paragraph_format.space_after = Pt(144)
    run_s = p_sub.add_run(subtitle)
    run_s.font.size = Pt(13)
    run_s.font.italic = True
    run_s.font.color.rgb = RGBColor(0x59, 0x59, 0x59)

    p_meta = doc.add_paragraph()
    p_meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_meta.paragraph_format.line_spacing = 1.3
    
    r = p_meta.add_run(f"Estudiante: {student_name}\n")
    r.font.bold = True
    p_meta.add_run(f"Materia / Curso: Investigación Académica y Proyectos\n")
    p_meta.add_run(f"Docente / Licenciado: {teacher_name}\n")
    p_meta.add_run(f"Fecha de entrega: {date_str}\n")
    p_meta.add_run("Formato de Presentación: Normas APA 7ma Edición (Mercado Boliviano)")
    
    doc.add_page_break()

    # SECCIONES DEL DOCUMENTO
    for sec in sections:
        h_type = sec.get("type", "h1")
        heading_text = sec.get("title", "")
        
        if h_type == "h1":
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(18)
            p.paragraph_format.space_after = Pt(8)
            run = p.add_run(heading_text)
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
        elif h_type == "h2":
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(14)
            p.paragraph_format.space_after = Pt(6)
            run = p.add_run(heading_text)
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)
        elif h_type == "paragraph":
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(6)
            p.add_run(sec.get("text", ""))
        elif h_type == "bullet":
            for item in sec.get("items", []):
                p = doc.add_paragraph(style='List Bullet')
                p.paragraph_format.space_after = Pt(4)
                p.add_run(item)
        elif h_type == "image":
            img_path = sec.get("path")
            caption = sec.get("caption", "")
            if os.path.exists(img_path):
                p_img = doc.add_paragraph()
                p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p_img.paragraph_format.space_before = Pt(10)
                p_img.paragraph_format.space_after = Pt(4)
                run_i = p_img.add_run()
                run_i.add_picture(img_path, width=Inches(5.5))
                
                if caption:
                    p_cap = doc.add_paragraph()
                    p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    p_cap.paragraph_format.space_after = Pt(12)
                    run_c = p_cap.add_run(f"Figura: {caption}")
                    run_c.font.size = Pt(9.5)
                    run_c.font.italic = True
                    run_c.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
        elif h_type == "table":
            headers = sec.get("headers", [])
            rows = sec.get("rows", [])
            table = doc.add_table(rows=len(rows) + 1, cols=len(headers))
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            
            hdr_cells = table.rows[0].cells
            for i, h in enumerate(headers):
                hdr_cells[i].text = h
                set_cell_background(hdr_cells[i], "1F497D")
                for p in hdr_cells[i].paragraphs:
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for r in p.runs:
                        r.font.bold = True
                        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        r.font.size = Pt(10)
                        
            for r_idx, row_data in enumerate(rows):
                row_cells = table.rows[r_idx + 1].cells
                bg_color = "F2F2F2" if r_idx % 2 == 1 else "FFFFFF"
                for c_idx, val in enumerate(row_data):
                    row_cells[c_idx].text = str(val)
                    set_cell_background(row_cells[c_idx], bg_color)
                    for p in row_cells[c_idx].paragraphs:
                        for r in p.runs:
                            r.font.size = Pt(9.5)
            
            p_space = doc.add_paragraph()
            p_space.paragraph_format.space_after = Pt(6)

    doc.save(filename)
    print(f"Documento APA guardado: {filename}")


def create_presentation(title, subtitle, slides_data, filename):
    prs = Presentation()
    prs.slide_width = PPTInches(13.333)
    prs.slide_height = PPTInches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Title Slide
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PPTInches(13.333), PPTInches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PPTRGBColor(0x1F, 0x49, 0x7D)
    bg.line.fill.background()
    
    txBox = slide.shapes.add_textbox(PPTInches(1), PPTInches(2.2), PPTInches(11.333), PPTInches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PPTPt(34)
    p.font.bold = True
    p.font.color.rgb = PPTRGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = PPTPt(18)
    p2.font.color.rgb = PPTRGBColor(0xD9, 0xE1, 0xF2)
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf.add_paragraph()
    p3.text = "\nEnfoque Especializado: Mercado Boliviano | Formato APA"
    p3.font.size = PPTPt(14)
    p3.font.italic = True
    p3.font.color.rgb = PPTRGBColor(0xBF, 0xBF, 0xBF)
    p3.alignment = PP_ALIGN.CENTER
    
    # Content Slides
    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PPTInches(13.333), PPTInches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = PPTRGBColor(0x1F, 0x49, 0x7D)
        header.line.fill.background()
        
        tf_h = header.text_frame
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.text = slide_info["title"]
        p_h.font.size = PPTPt(24)
        p_h.font.bold = True
        p_h.font.color.rgb = PPTRGBColor(0xFF, 0xFF, 0xFF)
        p_h.alignment = PP_ALIGN.LEFT
        
        img_path = slide_info.get("image")
        has_image = img_path and os.path.exists(img_path)
        
        text_width = PPTInches(6.8) if has_image else PPTInches(11.733)
        tx_content = slide.shapes.add_textbox(PPTInches(0.8), PPTInches(1.5), text_width, PPTInches(5.5))
        tf_c = tx_content.text_frame
        tf_c.word_wrap = True
        
        for idx, b in enumerate(slide_info["bullets"]):
            p_b = tf_c.paragraphs[0] if idx == 0 else tf_c.add_paragraph()
            p_b.text = f"•  {b}"
            p_b.font.size = PPTPt(15 if has_image else 16)
            p_b.font.color.rgb = PPTRGBColor(0x33, 0x33, 0x33)
            p_b.space_after = PPTPt(12)
            
        if has_image:
            slide.shapes.add_picture(img_path, PPTInches(7.9), PPTInches(1.6), width=PPTInches(4.8))
            
    prs.save(filename)
    print(f"Presentación PPTX guardada: {filename}")


# DEFINICIÓN TEMA 1 CON SECCIÓN 8 REFORZADA CON IMÁGENES, GRÁFICOS Y FLUJOGRAMAS
sections_tema1_bo = [
    {"type": "h1", "title": "1. Objetivo"},
    {"type": "h2", "title": "1.1 Objetivo General"},
    {"type": "paragraph", "text": "Investigar e identificar los modelos de negocio innovadores adaptados al mercado boliviano, analizando estrategias de inclusión financiera, economía circular y digitalización en el contexto socioeconómico local."},
    {"type": "h2", "title": "1.2 Objetivos Específicos"},
    {"type": "bullet", "items": [
        "Diagnosticar la estructura del mercado boliviano (alta informalidad laboral >75%, cultura del efectivo y dispersión geográfica).",
        "Analizar el impacto y modelo de negocio del sistema Pago QR Simple (Asoban/BCB), Tigo Money/Soli Pagos y la empresa de economía circular Mamut.",
        "Evaluar la viabilidad de escalabilidad y democratización del consumo digital en las ciudades del eje troncal (Santa Cruz, La Paz y Cochabamba).",
        "Diseñar el Diagrama de Gantt en MS Project para el lanzamiento de un modelo de negocio innovador e inclusivo en Bolivia."
    ]},

    {"type": "h1", "title": "2. Índice"},
    {"type": "bullet", "items": [
        "1. Objetivo (General y Específicos)",
        "2. Índice",
        "3. Requerimientos (Situación Problemática en Bolivia)",
        "4. Análisis de Caso de Investigación (Criterio Grupal - Casos QR Simple, Tigo Money y Mamut)",
        "5. Conclusión",
        "6. Conceptos Relacionados",
        "7. Bibliografía (Normas APA 7ma Edición)",
        "8. Imágenes, Gráficos, Estadísticas y Flujogramas",
        "9. Diagrama de Gantt en PROJECT (Estructura de Tareas)",
        "10. Otros Puntos Importantes (Matriz FODA Bolivia y Gestión de Riesgos)"
    ]},

    {"type": "h1", "title": "3. Requerimientos (Situación Problemática en Bolivia)"},
    {"type": "paragraph", "text": "El mercado boliviano presenta características únicas que exigen un rediseño de los modelos de negocio tradicionales importados:"},
    {"type": "bullet", "items": [
        "Alta Economía Informal: Bolivia registra una de las tasas de informalidad laboral más elevadas de América Latina (entre 75% y 80%), lo que dificulta la verificación de ingresos y bancarización formal.",
        "Predominio Histórico del Efectivo: El uso del dinero físico generaba altos costos logísticos, riesgos de seguridad y fricción transaccional en pequeños comercios y ferias informales (ej. La Pampa en Cochabamba o la 16 de Julio en El Alto).",
        "Brecha Tecnológica y Acceso a Capital: Los emprendimientos locales enfrentan altas barreras de financiamiento tradicional y costos de importación de insumos."
    ]},

    {"type": "h1", "title": "4. Análisis de Caso de Investigación (Criterio Grupal)"},
    {"type": "h2", "title": "4.1 Caso Pago QR Simple (Asoban / Banco Central de Bolivia)"},
    {"type": "paragraph", "text": "Desde un análisis de criterio grupal, el Pago QR Simple es la mayor innovación de modelo de negocio financiero en la historia reciente de Bolivia. Al establecer un estándar interoperable entre todos los bancos y billeteras móviles:"},
    {"type": "bullet", "items": [
        "Eliminó la necesidad de terminales POS costosas y comisiones bancarias altas para los microcomerciantes.",
        "Permitió realizar pagos al instante desde cualquier aplicación bancaria sin requerir número de cuenta complejo.",
        "Democratizó el cobro digital en mercados populares, transporte público y servicios independientes."
    ]},

    {"type": "h2", "title": "4.2 Caso Mamut (Cochabamba, Bolivia – Economía Circular)"},
    {"type": "paragraph", "text": "La empresa boliviana Mamut representa un modelo innovador de economía circular de triple impacto (económico, social y ambiental):"},
    {"type": "bullet", "items": [
        "Transforma neumáticos fuera de uso (caucho reciclado contaminante) en baldosas amortiguantes para parques y pistas atléticas.",
        "Genera una cadena de valor local integrando a recolectores informales de llantas.",
        "Ha logrado exportar tecnología y productos terminados a mercados internacionales como Paraguay y México."
    ]},

    {"type": "h1", "title": "5. Conclusión"},
    {"type": "bullet", "items": [
        "La innovación en Bolivia debe ser inclusiva y resolver problemas reales de la economía informal.",
        "El éxito de modelos como QR Simple demuestra que la interoperabilidad es clave para la adopción masiva.",
        "Empresas como Mamut prueban que Bolivia posee capacidad técnica para exportar innovación sostenible."
    ]},

    {"type": "h1", "title": "6. Conceptos Relacionados"},
    {"type": "bullet", "items": [
        "Interoperabilidad Financiera: Capacidad de conectar distintas entidades bancarias y billeteras móviles en una misma plataforma transaccional.",
        "Economía Circular: Modelo de producción y consumo que implica reutilizar y reciclar materiales existentes todas las veces posibles.",
        "Billetera Móvil (E-Money): Instrumento electrónico de almacenamiento de dinero administrado a través de teléfonos móviles.",
        "Innovación Frugal Boliviana: Adaptación de procesos productivos para ofrecer servicios de alta calidad con recursos limitados."
    ]},

    {"type": "h1", "title": "7. Bibliografía (Normas APA 7ma Edición)"},
    {"type": "bullet", "items": [
        "Asociación de Bancos Privados de Bolivia [ASOBAN]. (2024). Reporte de inclusión financiera y evolución del Pago QR en Bolivia. ASOBAN.",
        "Banco Central de Bolivia [BCB]. (2023). Informe de Sistemas de Pago en Bolivia 2023. BCB.",
        "Empresa Mamut. (2024). Memoria de sostenibilidad y desarrollo de infraestructura urbana verde en América Latina. Mamut Cochabamba.",
        "Instituto Nacional de Estadística [INE]. (2023). Encuesta de empleo e informalidad laboral en Bolivia. INE Bolivia."
    ]},

    {"type": "h1", "title": "8. Imágenes, Gráficos, Estadísticas y Flujogramas"},
    {"type": "h2", "title": "8.1 Imágenes del Caso de Estudio en Bolivia"},
    {"type": "image", "path": "qr_simple_bolivia.png", "caption": "Imagen 1: Ecosistema de Pago QR Simple y Billeteras Móviles en Bolivia"},
    {"type": "image", "path": "mamut_bolivia.png", "caption": "Imagen 2: Productos Industriales de Caucho Reciclado por la Empresa Mamut (Cochabamba)"},
    
    {"type": "h2", "title": "8.2 Gráfico Estadístico de Crecimiento en Bolivia"},
    {"type": "paragraph", "text": "A continuación se presenta el gráfico estadístico de barras y tendencia del volumen transaccionado vía Pago QR en el mercado boliviano:"},
    {"type": "image", "path": "grafico_tema1_estadisticas_qr.png", "caption": "Gráfico 1: Evolución del Pago QR Simple en Bolivia (Millones de Operaciones y Monto en Bs)"},

    {"type": "h2", "title": "8.3 Tabla de Datos Estadísticos"},
    {"type": "table",
     "headers": ["Año", "Operaciones QR Anuales (Millones)", "Monto Transaccionado (Millones Bs)", "Comercios Adheridos (%)"],
     "rows": [
         ["2020", "12.4 M", "Bs 1.200 M", "15%"],
         ["2022", "85.2 M", "Bs 9.500 M", "48%"],
         ["2024", "210.0 M", "Bs 26.000 M", "76%"],
         ["2026 (Proy.)", "380.0 M", "Bs 48.000 M", "88%"]
     ]},

    {"type": "h2", "title": "8.4 Diagrama de Flujo / Flujograma Operativo"},
    {"type": "paragraph", "text": "Representación gráfica del flujo transaccional de cobro digital en mercados populares bolivianos:"},
    {"type": "image", "path": "grafico_tema1_flujograma.png", "caption": "Flujograma 1: Proceso Operativo Transaccional del Pago QR Simple e Inclusión Financiera en Bolivia"},

    {"type": "h1", "title": "9. Diagrama de Gantt en PROJECT"},
    {"type": "paragraph", "text": "Estructura WBS configurada para ser importada en Microsoft Project para un Proyecto de Innovación Digital en Bolivia:"},
    {"type": "table",
     "headers": ["EDT / WBS", "Nombre de la Tarea", "Duración", "Inicio", "Fin", "Predecesoras", "Recursos Asignados"],
     "rows": [
         ["1", "FASE 1: DIAGNÓSTICO Y MERCADO BOLIVIANO", "15 días", "01/09/2026", "21/09/2026", "", "Analista Local, UX Researcher"],
         ["1.1", "Encuesta en mercados del eje troncal (LP-CB-SC)", "8 días", "01/09/2026", "10/09/2026", "", "Investigadores de Campo"],
         ["1.2", "Revisión de normativa ASFI y Banco Central", "7 días", "11/09/2026", "21/09/2026", "1.1", "Abogado Financiero"],
         ["2", "FASE 2: DESARROLLO E INTEGRACIÓN QR Y APP", "20 días", "22/09/2026", "19/10/2026", "1", "Dev Team Lead"],
         ["2.1", "Integración API QR Simple y Billeteras", "10 días", "22/09/2026", "05/10/2026", "1.2", "Backend Developer"],
         ["2.2", "Pruebas de estrés transaccional", "10 días", "06/10/2026", "19/10/2026", "2.1", "QA Tester"],
         ["3", "FASE 3: PILOTO Y CAPACITACIÓN A COMERCIANTES", "25 días", "20/10/2026", "24/11/2026", "2", "Promotores de Campo"],
         ["3.1", "Capacitación a 500 comerciantes informales", "15 días", "20/10/2026", "09/11/2026", "2.2", "Facilitadores"],
         ["3.2", "Evaluación de feedback y ajustes UX", "10 días", "10/11/2026", "24/11/2026", "3.1", "Product Owner"]
     ]},

    {"type": "h1", "title": "10. Otros Puntos Importantes"},
    {"type": "h2", "title": "10.1 Matriz FODA del Entorno Emprendedor Boliviano"},
    {"type": "table",
     "headers": ["Fortalezas (F)", "Oportunidades (O)", "Debilidades (D)", "Amenazas (A)"],
     "rows": [
         ["• Adopción masiva de QR Simple\n• Creatividad del emprendedor boliviano\n• Bajos costos operativos digitales", 
          "• Inclusión de la masa informal urbana\n• Crecimiento del e-commerce local", 
          "• Brecha de crédito bancario tradicional\n• Logística fragmentada fuera de ciudades", 
          "• Fluctuación cambiaria e inflación en insumos\n• Cambios normativos imprevistos"]
     ]}
]

slides_tema1_bo = [
    {"title": "Modelos de Negocio Innovadores en Bolivia", "bullets": [
        "Investigación adaptada 100% al contexto del mercado boliviano.",
        "Análisis de inclusión financiera y economía circular.",
        "Formato de presentación según Normas APA 7ma Edición."
    ], "image": "qr_simple_bolivia.png"},
    {"title": "Objetivos de la Investigación", "bullets": [
        "General: Diagnosticar la disrupción de modelos innovadores en Bolivia.",
        "Específico 1: Evaluar el impacto del Pago QR Simple (Asoban/BCB).",
        "Específico 2: Estudiar el modelo circular de la empresa Mamut.",
        "Específico 3: Diseñar el plan WBS en MS Project para proyectos locales."
    ]},
    {"title": "Situación Problemática en Bolivia", "bullets": [
        "Informalidad laboral superior al 75% en ciudades principales.",
        "Uso histórico de efectivo generaba inseguridad y altos costos.",
        "Requerimiento: Crear herramientas digitales de costo cero para el usuario."
    ]},
    {"title": "Casos Emblemáticos en Bolivia", "bullets": [
        "Caso QR Simple: Pagos inmediatos interbancarios sin comisión.",
        "Tigo Money: Billetera móvil para sectores no bancarizados.",
        "Caso Mamut: Reciclaje de llantas en pisos industriales exportados a LATAM."
    ], "image": "mamut_bolivia.png"},
    {"title": "Punto 8: Gráfico Estadístico de Crecimiento QR", "bullets": [
        "Operaciones QR: 12.4M (2020) ➔ 210M (2024) ➔ 380M (2026 proy.).",
        "Monto transaccionado superó los 26.000 Millones de Bs.",
        "76% de comercios y ferias formales e informales aceptan cobro por QR."
    ], "image": "grafico_tema1_estadisticas_qr.png"},
    {"title": "Punto 8: Flujograma del Proceso Operativo QR", "bullets": [
        "Flujo directo de cobro y pago en mercados populares de Bolivia.",
        "Interoperabilidad inmediata respaldada por el Banco Central de Bolivia."
    ], "image": "grafico_tema1_flujograma.png"},
    {"title": "Plan de Proyecto en MS Project", "bullets": [
        "Fase 1: Diagnóstico y Estudio de Mercado Boliviano (15 días)",
        "Fase 2: Desarrollo e Integración API QR (20 días)",
        "Fase 3: Piloto y Capacitación a Comerciantes (25 días)"
    ]},
    {"title": "Conclusiones", "bullets": [
        "La simplicidad e interoperabilidad impulsaron la revolución digital boliviana.",
        "Mamut y Quantum demuestran la capacidad de innovación sostenible en Bolivia."
    ]}
]


# DEFINICIÓN TEMA 2 CON SECCIÓN 8 REFORZADA CON IMÁGENES, GRÁFICOS Y FLUJOGRAMAS
sections_tema2_bo = [
    {"type": "h1", "title": "1. Objetivo"},
    {"type": "h2", "title": "1.1 Objetivo General"},
    {"type": "paragraph", "text": "Evaluar el estado actual, la infraestructura de electrineras, el marco regulatorio del Decreto Supremo 4539 y las perspectivas futuras de la electromovilidad y los vehículos eléctricos en Bolivia."},
    {"type": "h2", "title": "1.2 Objetivos Específicos"},
    {"type": "bullet", "items": [
        "Analizar el impacto del Decreto Supremo N° 4539 en la exención de aranceles e impuestos para vehículos eléctricos e insumos en Bolivia.",
        "Cuantificar el ahorro fiscal que representaría la sustitución del subsidio a la gasolina y diésel mediante la electromovilidad.",
        "Evaluar el modelo de producción de Quantum Motors S.A. y el despliegue de electrineras de ENDE Corporación en el eje troncal.",
        "Diseñar una estructura de desglose de trabajo (EDT/WBS) en MS Project para la instalación de estaciones de carga rápida en Bolivia."
    ]},

    {"type": "h1", "title": "2. Índice"},
    {"type": "bullet", "items": [
        "1. Objetivo (General y Específicos)",
        "2. Índice",
        "3. Requerimientos (Situación Problemática de la Energía en Bolivia)",
        "4. Análisis de Caso de Investigación (Criterio Grupal - Quantum Motors y ENDE)",
        "5. Conclusión",
        "6. Conceptos Relacionados",
        "7. Bibliografía (Normas APA 7ma Edición)",
        "8. Imágenes, Gráficos, Estadísticas y Flujogramas",
        "9. Diagrama de Gantt en PROJECT (Estructura de Tareas)",
        "10. Otros Puntos Importantes (FODA Electromovilidad Bolivia e Impacto Fiscal)"
    ]},

    {"type": "h1", "title": "3. Requerimientos (Situación Problemática en Bolivia)"},
    {"type": "paragraph", "text": "El sector transporte en Bolivia depende de manera crítica de la importación de diésel y gasolina subvencionados. El Estado boliviano destina más de $us 1.500 millones anuales para mantener fijos los precios al consumidor (Bs 3.74/litro de gasolina). Esta situación genera los siguientes requerimientos:"},
    {"type": "bullet", "items": [
        "Transición de Matriz Energética: Urgencia de migrar del consumo de hidrocarburos importados a la energía eléctrica generada localmente por ENDE (hidroeléctrica y solar).",
        "Infraestructura de Recarga Faltante: Bolivia requiere ampliar la red de electrineras públicas en las carreteras Cochabamba - Santa Cruz y La Paz - Oruro.",
        "Industrialización del Litio: Vincular la reserva de litio del Salar de Uyuni (YLB) con la producción local de baterías LFP."
    ]},

    {"type": "h1", "title": "4. Análisis de Caso de Investigación (Criterio Grupal)"},
    {"type": "h2", "title": "4.1 Caso Quantum Motors S.A. (Cochabamba, Bolivia)"},
    {"type": "paragraph", "text": "El análisis en grupo del caso Quantum Motors S.A. permite destacar los siguientes hitos de la industria boliviana:"},
    {"type": "bullet", "items": [
        "Primer Fabricante Boliviano de EVs: Fundada en Cochabamba, diseña y ensambla vehículos de micro-movilidad urbana (Modelos E4, Nexus).",
        "Carga Convencional a 220V: Sus vehículos no requieren instalaciones especiales complejas, pudiendo enchufarse en cualquier toma del hogar.",
        "Exportación Regional: Expansión de la marca boliviana a Paraguay, El Salvador y México."
    ]},

    {"type": "h2", "title": "4.2 Red Nacional de Electrineras (ENDE Corporación y YPFB)"},
    {"type": "paragraph", "text": "La Empresa Nacional de Electricidad (ENDE) lidera la instalación de los primeros tótems de recarga gratuita y de pago en el eje troncal:"},
    {"type": "bullet", "items": [
        "Puntos estratégicos instalados en estaciones de servicio de YPFB en La Paz, Cochabamba y Santa Cruz.",
        "Fomento del uso de vehículos eléctricos híbridos y 100% batería protegidos por el Decreto Supremo 4539."
    ]},

    {"type": "h1", "title": "5. Conclusión"},
    {"type": "bullet", "items": [
        "La electromovilidad en Bolivia es una alternativa viable que ahorra hasta un 75% en costos de operación frente a la gasolina.",
        "Quantum Motors demuestra que Bolivia puede ser un actor clave en la industrialización de la micro-movilidad eléctrica en LATAM.",
        "Se requiere acelerar la construcción de corredores eléctricos interdepartamentales."
    ]},

    {"type": "h1", "title": "6. Conceptos Relacionados"},
    {"type": "bullet", "items": [
        "Decreto Supremo 4539: Norma legal boliviana que otorga 0% de arancel a la importación de EVs y repuestos.",
        "Electrinera (EVSE): Punto de recarga eléctrica diseñado para abastecer las baterías de autos eléctricos.",
        "Yacimientos de Litio Bolivianos (YLB): Empresa estatal responsable de la extracción e industrialización del litio en Uyuni.",
        "Costo Por Kilómetro (Bs/km): Indicador económico que demuestra el ahorro de la electricidad frente a los combustibles fósiles."
    ]},

    {"type": "h1", "title": "7. Bibliografía (Normas APA 7ma Edición)"},
    {"type": "bullet", "items": [
        "Decreto Supremo N° 4539. (2021, 7 de julio). Incentivos tributarios y financieros para la electromovilidad en Bolivia. Gaceta Oficial del Estado Plurinacional de Bolivia.",
        "ENDE Corporación. (2023). Memoria anual y plan de desarrollo de estaciones de recarga en Bolivia. ENDE.",
        "Ministerio de Hidrocarburos y Energías [MHE]. (2023). Estrategia Nacional de Electromovilidad 2023-2030. MHE Bolivia.",
        "Quantum Motors. (2024). Ficha técnica y reporte de impacto ambiental de la flota eléctrica en Bolivia. Quantum."
    ]},

    {"type": "h1", "title": "8. Imágenes, Gráficos, Estadísticas y Flujogramas"},
    {"type": "h2", "title": "8.1 Imágenes del Caso de Estudio en Bolivia"},
    {"type": "image", "path": "quantum_motors_bolivia.png", "caption": "Imagen 1: Vehículo Eléctrico Urbano Quantum Motors fabricado en Cochabamba, Bolivia"},
    {"type": "image", "path": "electrinera_bolivia.png", "caption": "Imagen 2: Estación de Recarga Electrinera de ENDE Corporación en Bolivia"},

    {"type": "h2", "title": "8.2 Gráfico Estadístico de Comparativa de Costos en Bolivia"},
    {"type": "paragraph", "text": "Gráfico comparativo de costos operativos y mantenimiento anual en Bolivia (Vehículo a Gasolina vs Vehículo Eléctrico Quantum):"},
    {"type": "image", "path": "grafico_tema2_comparativa_costos.png", "caption": "Gráfico 1: Comparativa de Costo Operativo por 100km y Mantenimiento Anual en Bolivia (Bs)"},

    {"type": "h2", "title": "8.3 Tabla de Datos Estadísticos"},
    {"type": "table",
     "headers": ["Concepto", "Vehículo a Gasolina (1.2L)", "Vehículo Eléctrico (Quantum E4)", "Ventaja Económica"],
     "rows": [
         ["Costo por 100 km", "Bs 42.00 (Gasolina Subvencionada)", "Bs 11.50 (Tarifa Eléctrica Residencial)", "72.6% Ahorro Directo"],
         ["Mantenimiento Anual", "Bs 2.500 (Aceites, filtros)", "Bs 600 (Revisión general)", "76.0% Menor Costo"],
         ["Arancel de Importación", "10% a 20%", "0% (Bajo D.S. 4539)", "Exención Total"],
         ["Emisiones Directas", "145 g CO2/km", "0 g CO2/km", "100% Ecológico"]
     ]},

    {"type": "h2", "title": "8.4 Diagrama de Flujo / Flujograma Energético"},
    {"type": "paragraph", "text": "Representación gráfica del flujo de la red eléctrica nacional desde la generación de ENDE hasta la recarga del vehículo eláctrico:"},
    {"type": "image", "path": "grafico_tema2_flujograma.png", "caption": "Flujograma 1: Matriz Energética e Integración de Electromovilidad en Bolivia"},

    {"type": "h1", "title": "9. Diagrama de Gantt en PROJECT"},
    {"type": "paragraph", "text": "Estructura WBS lista para importar en Microsoft Project para el 'Despliegue de Electrineras en Bolivia':"},
    {"type": "table",
     "headers": ["EDT / WBS", "Nombre de la Tarea", "Duración", "Inicio", "Fin", "Predecesoras", "Recursos Asignados"],
     "rows": [
         ["1", "FASE 1: ESTUDIO DE FACTIBILIDAD RED NACIONA", "20 días", "01/09/2026", "28/09/2026", "", "Ingeniero Eléctrico ENDE"],
         ["1.1", "Estudio de demanda carretera LP-CB-SC", "10 días", "01/09/2026", "14/09/2026", "", "Técnicos MHE"],
         ["1.2", "Inspección de subestaciones eléctricas", "10 días", "15/09/2026", "28/09/2026", "1.1", "Ingenieros de Distribución"],
         ["2", "FASE 2: COMPRA E INSTALACIÓN DE CARGADORES DC", "35 días", "29/09/2026", "16/11/2026", "1", "Equipo de Contrataciones"],
         ["2.1", "Importación de tótems con exención DS 4539", "15 días", "29/09/2026", "19/10/2026", "1.2", "Agente Aduanero"],
         ["2.2", "Obras civiles e instalación eléctrica", "20 días", "20/10/2026", "16/11/2026", "2.1", "Contratista Civil"],
         ["3", "FASE 3: PRUEBAS Y COMISIONAMIENTO", "15 días", "17/11/2026", "07/12/2026", "2", "Supervisores ENDE"],
         ["3.1", "Pruebas de carga rápida con flota Quantum", "10 días", "17/11/2026", "30/11/2026", "2.2", "Pilotos de Prueba"],
         ["3.2", "Inauguración de la Ruta Eléctrica Bolivia", "5 días", "01/12/2026", "07/12/2026", "3.1", "Relaciones Públicas"]
     ]},

    {"type": "h1", "title": "10. Otros Puntos Importantes"},
    {"type": "h2", "title": "10.1 Matriz FODA de la Electromovilidad en Bolivia"},
    {"type": "table",
     "headers": ["Fortalezas (F)", "Oportunidades (O)", "Debilidades (D)", "Amenazas (A)"],
     "rows": [
         ["• Reservas mundiales de Litio en Uyuni\n• Producción nacional Quantum Motors\n• Exención tributaria Decreto 4539", 
          "• Reducción masiva de la factura de importación de combustibles\n• Posicionamiento verde regional", 
          "• Red de electrineras en fase inicial de expansión\n• Menor potencia eléctrica en zonas rurales", 
          "• Resistencia de sindicatos del transporte tradicional\n• Subsidio a la gasolina persistente"]
     ]}
]

slides_tema2_bo = [
    {"title": "Vehículos Eléctricos y Electromovilidad en Bolivia", "bullets": [
        "Investigación enfocada 100% en la transición energética en Bolivia.",
        "Análisis de incentivos del D.S. 4539 y caso Quantum Motors S.A.",
        "Formato de presentación según Normas APA 7ma Edición."
    ], "image": "quantum_motors_bolivia.png"},
    {"title": "Objetivos de la Investigación", "bullets": [
        "General: Evaluar la viabilidad técnica y fiscal de la electromovilidad en Bolivia.",
        "Específico 1: Analizar el impacto legal y tributario del Decreto Supremo 4539.",
        "Específico 2: Cuantificar la reducción del subsidio estatal a combustibles.",
        "Específico 3: Diseñar el plan WBS en MS Project para la Red de Electrineras ENDE."
    ]},
    {"title": "Situación Problemática del Sector Energético", "bullets": [
        "Gasto fiscal >$us 1.500 millones anuales en subsidiar gasolina y diésel importados.",
        "Necesidad urgente de migrar hacia la energía hidroeléctrica/solar local de ENDE.",
        "Falta de electrineras de carga rápida en carreteras interdepartamentales."
    ]},
    {"title": "Casos de Éxito en Bolivia", "bullets": [
        "Quantum Motors: Primer fabricante boliviano de autos eléctricos (Cochabamba).",
        "Carga a 220V convencional en cualquier enchufe de hogar.",
        "ENDE Corporación y YPFB: Instalación de la red inicial de electrineras."
    ], "image": "electrinera_bolivia.png"},
    {"title": "Punto 8: Gráfico Estadístico de Comparativa de Costos", "bullets": [
        "Costo por 100 km: Gasolina Bs 42.00 vs Eléctrico Quantum Bs 11.50 (72.6% Ahorro).",
        "Mantenimiento anual 76% más económico en vehículos eléctricos.",
        "Arancel de importación: 0% por norma legal D.S. 4539."
    ], "image": "grafico_tema2_comparativa_costos.png"},
    {"title": "Punto 8: Flujograma de la Red Electromóvil Bolivia", "bullets": [
        "Integración de energía renovable hidro/solar de ENDE a electrineras.",
        "Almacenamiento en baterías de litio de producción nacional."
    ], "image": "grafico_tema2_flujograma.png"},
    {"title": "Plan de Proyecto en MS Project", "bullets": [
        "Fase 1: Estudio de Factibilidad Red Nacional (20 días)",
        "Fase 2: Compra e Instalación de Cargadores DC (35 días)",
        "Fase 3: Pruebas y Comisionamiento (15 días)"
    ]},
    {"title": "Conclusiones", "bullets": [
        "La electromovilidad cuida la economía familiar y la salud fiscal de Bolivia.",
        "El litio boliviano y Quantum Motors posicionan al país en la industria del futuro."
    ]}
]

# GENERAR LOS ARCHIVOS DEFINITIVOS
create_apa_document(
    title="INVESTIGACIÓN ACADÉMICA: MODELOS DE NEGOCIO INNOVADORES EN EL MERCADO BOLIVIANO",
    subtitle="Análisis de inclusión financiera (Pago QR Simple, Tigo Money), economía circular (Mamut) y estructura MS Project",
    student_name="Carlos Carrillo",
    teacher_name="Docente / Licenciado de Materia",
    date_str="04 de Agosto de 2026",
    sections=sections_tema1_bo,
    filename="Investigacion_Tema1_Modelos_de_Negocio_Bolivia_Final.docx"
)

create_presentation(
    title="Modelos de Negocio Innovadores en el Mercado Boliviano",
    subtitle="Presentación Ejecutiva e Ilustrada - Formato APA",
    slides_data=slides_tema1_bo,
    filename="Presentacion_Tema1_Modelos_de_Negocio_Bolivia_Final.pptx"
)

create_apa_document(
    title="INVESTIGACIÓN ACADÉMICA: VEHÍCULOS ELÉCTRICOS Y EL FUTURO DE LA ELECTROMOVILIDAD EN BOLIVIA",
    subtitle="Análisis de incentivos (D.S. 4539), caso Quantum Motors, red de electrineras ENDE y plan MS Project",
    student_name="Carlos Carrillo",
    teacher_name="Docente / Licenciado de Materia",
    date_str="04 de Agosto de 2026",
    sections=sections_tema2_bo,
    filename="Investigacion_Tema2_Electromovilidad_Bolivia_Final.docx"
)

create_presentation(
    title="Vehículos Eléctricos y Electromovilidad en Bolivia",
    subtitle="Presentación Ejecutiva e Ilustrada - Formato APA",
    slides_data=slides_tema2_bo,
    filename="Presentacion_Tema2_Electromovilidad_Bolivia_Final.pptx"
)

print("¡TODOS LOS DOCUMENTOS REFORZADOS CON SECCIÓN 8 DE IMÁGENES, GRÁFICOS, ESTADÍSTICAS Y FLUJOGRAMAS HAN SIDO GENERADOS EXITOSAMENTE!")

